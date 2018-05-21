# -*- coding: UTF-8 -*-
''' combine sentences in different notes into paragraphs '''
import os
import csv
from pptx import Presentation
import re
import numpy as np # average, standard deviation, max
from nltk.tokenize import sent_tokenize
import sys

SUBJECT = 'OM'
CHAPTER = 'supplementA'

SUBJECT = sys.argv[1]
CHAPTER = sys.argv[2]
print (SUBJECT, CHAPTER)

NOTE_FOLDER = 'note/' + SUBJECT + '/' + CHAPTER + '/'
SLIDE_FOLDER = 'ppt/' + SUBJECT + '/'
TOPIC_FOLDER = 'topic/' + SUBJECT + '/'
SLIDE_FILE_NAME = CHAPTER + '.pptx'
COMBINEDSLIDE_FILE_NAME = CHAPTER + '_slides.csv'
TOPIC_FILE_NAME = 'Topics_' + CHAPTER + '.csv'
MIXEDNOTEPARA_FILE = 'MixedNote_' + CHAPTER + '.csv'

Topics = []
with open (TOPIC_FOLDER + TOPIC_FILE_NAME, 'r', encoding = 'utf-8') as termFile:
    reader = csv.DictReader(termFile)
    for row in reader:
        Topics.append({'term_e': row['term_e'], 'abbr': row['abbr'], 'term_c': row['term_c']})

with open ('stopwords.txt', 'r', encoding = 'utf-8') as inFile:
    StopWords = [sw.replace('\n', '') for sw in inFile.readlines()]

RawNoteWordsList = []
PunctuationMarks = ['.', '!', '?', ':', ';', '\'', '\"', ',', '，', '。', '！', '？', '：', '；', '「', '」', '『', '』', '※']
SpecialMarks = ['+', '-', '*', '/', '=', '\\', '<', '>', '~', '@', '#', '(', ')', '[', ']', '{', '}', '|', '^', '–', '—']

prs = Presentation(SLIDE_FOLDER + SLIDE_FILE_NAME)
slides = prs.slides
Slides = []
text_runs = []
sid = 0
for slide in prs.slides:
    title = ''
    if slide.shapes.title:
        title = slide.shapes.title.text
    for shape in slide.shapes:
        content = ''
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            text = ''
            for run in paragraph.runs:
                if len(run.text) > 0: text += run.text
            if len(text) > 0:
                # text.replace('\n', ' ')
                # if re.match('[A-Z][a-z].+', text) or re.match('\d', text):
                #   text = '\n' + text
                # content += text.replace('\t', ' ') + ' '
                content += text.replace('\n', ' ').replace('\t', ' ') + ' '
        content = re.sub(' +', ' ', content) # remove duplicate spaces
        content = '\n'.join(sent_tokenize(content)) # split sentences
        if sid > 0 and title == Slides[sid - 1]['title']:
            Slides[sid - 1]['content'] = Slides[sid - 1]['content'] + '\n' + content.strip() # combine the slides which have a same title
        else:
            Slides.append({'title': title, 'content': content.strip()})
            sid += 1
for sid in range(len(Slides)):
    Slides[sid]['topic'] = []
    for term in Topics:
        if term['term_e'] in Slides[sid]['content'] or (term['abbr'] != '' and (' ' + term['abbr'] + ' ' in Slides[sid]['content'] or '(' + term['abbr'] + ')' in Slides[sid]['content'])):
            Slides[sid]['topic'].append(term['term_e'])

with open (SLIDE_FOLDER + COMBINEDSLIDE_FILE_NAME, 'w', newline = '', encoding = 'utf-8') as outFile:
    writer = csv.DictWriter(outFile, fieldnames = ['title', 'content', 'topic'])
    writer.writeheader()
    for slide in Slides:
        writer.writerow(slide)