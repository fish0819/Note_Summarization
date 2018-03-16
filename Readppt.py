''' read ppt '''
from pptx import Presentation
import re

FILE_NAME = 'ppt\\Supplement-A.pptx'
prs = Presentation(FILE_NAME)
slides = prs.slides
NoteList = []
Slides = []
text_runs = []
sid = 0
for slide in prs.slides:
	title = ''
	if slide.has_notes_slide:
		for shape in slide.notes_slide.shapes:
			if shape.text != '':
				NoteList.append(shape.text)
	if slide.shapes.title:
		title = slide.shapes.title.text
	for shape in slide.shapes:
		content = ''
		if not shape.has_text_frame:
			continue
		for paragraph in shape.text_frame.paragraphs:
			text = ''
			for run in paragraph.runs:
				if len(run.text) > 0: text += run.text
			if len(text) > 0:
				content += text.replace('\n', ' ').replace('\t', ' ') + ' '
		content = re.sub(' +', ' ', content) # remove duplicate spaces
		if sid > 0 and title == Slides[sid - 1]['title']:
			Slides[sid - 1]['content'] = Slides[sid - 1]['content'] + ' ' + content.strip() # combine the slides which have a same title
		else:
			Slides.append({'title': title, 'content': content.strip()})
			sid += 1
for s in Slides:
	print (s)
	print ()